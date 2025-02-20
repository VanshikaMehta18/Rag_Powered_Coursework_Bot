import os
from flask import Flask, request, jsonify, render_template
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains import RetrievalQA
from langchain_openai import ChatOpenAI
from dotenv import load_dotenv
import openai
from langchain_community.embeddings import SentenceTransformerEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings

app = Flask(__name__)

load_dotenv("path-to-your-.env-file")

openai_api_key = os.getenv("OPENAI_API_KEY")
openai.api_key = openai_api_key

def extract_text_from_pptx(file_path):
    presentation = Presentation(file_path)
    text = []
    
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

        if slide_text:
            text.append("\n".join(slide_text))

    return "\n\n".join(text) if text else None

ppt_folder = r"path-to-the-folder-which-contains-all-ppts"
documents = []

for file in os.listdir(ppt_folder):
    file_path = os.path.join(ppt_folder, file)

    if file.endswith(".pptx"):
        text = extract_text_from_pptx(file_path)

        if text:
            documents.append(text)

if not documents:
    raise ValueError("No valid text was extracted from the PPTX files.")

course_text = "\n\n".join(documents)

text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=100)
chunks = text_splitter.split_text(course_text)

if not chunks:
    raise ValueError("No chunks were created from the extracted text.")

embedding_model = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
vector_store = FAISS.from_texts(chunks, embedding_model)

llm = ChatOpenAI(model_name="gpt-4o-mini", openai_api_key=openai_api_key)

retriever = vector_store.as_retriever()

qa_chain = RetrievalQA.from_chain_type(llm=llm, retriever=retriever)

@app.route("/")
def home():
    return render_template("index.html")

@app.route("/chat", methods=["GET"])
def chat():
    user_query = request.args.get("query")
    
    if user_query:
        try:
            response = qa_chain.invoke({"query": user_query})
            result = response.get("result", "I couldn't understand that.")
            return jsonify({"response": result})
        except Exception as e:
            return jsonify({"error": str(e)})
    else:
        return jsonify({"error": "No query provided."})

if __name__ == "__main__":
    app.run(debug=True)
